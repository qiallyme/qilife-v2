#!/usr/bin/env python3
# src/fileflow/content_extractor.py

import sys
import os
import subprocess
from pathlib import Path
from datetime import datetime

# Text extraction
import pytesseract            # pip install pytesseract
from PIL import Image         # pip install Pillow
import pdfplumber             # pip install pdfplumber
import docx                   # pip install python-docx
from pptx import Presentation # pip install python-pptx
import pandas as pd           # pip install pandas

# Audio/video transcription (lazy OpenAI client)
from openai import OpenAI     # pip install openai

def extract_text(path: Path) -> str:
    suffix = path.suffix.lower()

    # 1. Plain text & code
    if suffix in {".txt", ".md", ".py", ".js", ".java", ".html", ".css", ".json", ".xml", ".eml", ".mhtml"}:
        return path.read_text(encoding="utf-8", errors="ignore")

    # 2. PDF
    if suffix == ".pdf":
        pages = []
        with pdfplumber.open(path) as pdf:
            for p in pdf.pages:
                pages.append(p.extract_text() or "")
        return "\n".join(pages)

    # 3. Word .docx
    if suffix == ".docx":
        doc = docx.Document(path)
        return "\n".join(p.text for p in doc.paragraphs)

    # 4. PowerPoint .pptx
    if suffix == ".pptx":
        prs = Presentation(str(path))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
        return "\n".join(texts)

    # 5. Spreadsheets (.xls/.xlsx/.csv)
    if suffix in {".xls", ".xlsx"}:
        try:
            df = pd.read_excel(path, engine="openpyxl")
        except Exception:
            df = pd.read_csv(path)
        return df.to_csv(index=False)
    if suffix == ".csv":
        df = pd.read_csv(path)
        return df.to_csv(index=False)

    # 6. Images → OCR
    if suffix in {".png", ".jpg", ".jpeg", ".tiff", ".bmp", ".gif"}:
        img = Image.open(path)
        return pytesseract.image_to_string(img)

    # 7. Audio → Whisper
    if suffix in {".wav", ".mp3", ".m4a", ".flac"}:
        return _transcribe_audio(path)

    # 8. Video → extract audio + Whisper
    if suffix in {".mp4", ".mov", ".avi", ".mkv"}:
        return _transcribe_audio(path, is_video=True)

    # 9. Fallback → nothing
    return ""

def _transcribe_audio(path: Path, is_video: bool = False) -> str:
    """
    Uses ffmpeg to convert/extract audio and Whisper for transcription.
    Requires OPENAI_API_KEY in environment.
    """
    api_key = os.getenv("OPENAI_API_KEY")
    if not api_key:
        print("[yellow]⚠️ OPENAI_API_KEY not set – skipping transcription[/]")
        return ""

    # prepare wav
    wav = path.with_suffix(".wav")
    cmd = ["ffmpeg", "-y", "-i", str(path)]
    if is_video:
        cmd += ["-vn"]
    cmd += ["-ar", "16000", "-ac", "1", str(wav)]
    subprocess.run(cmd, check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)

    # lazy-init client
    client = OpenAI(api_key=api_key)
    with open(wav, "rb") as f:
        resp = client.audio.transcriptions.create(model="whisper-1", file=f)
    return resp.text

def extract_metadata(path: Path) -> dict:
    stat = path.stat()
    return {
        "filename": path.name,
        "extension": path.suffix.lower(),
        "size_bytes": stat.st_size,
        "created": datetime.fromtimestamp(stat.st_ctime).isoformat(),
        "modified": datetime.fromtimestamp(stat.st_mtime).isoformat(),
    }

def extract_context(path: Path) -> dict:
    """
    Returns:
      - 'text': extracted text or transcript
      - 'metadata': file metadata
    """
    return {
        "text": extract_text(path),
        "metadata": extract_metadata(path)
    }

if __name__ == "__main__":
    if len(sys.argv) != 2:
        print("Usage: content_extractor.py /path/to/file")
        sys.exit(1)
    p = Path(sys.argv[1])
    if not p.exists():
        print(f"File not found: {p}")
        sys.exit(1)
    ctx = extract_context(p)
    print("--- METADATA ---")
    for k, v in ctx["metadata"].items():
        print(f"{k}: {v}")
    print("\n--- TEXT PREVIEW ---\n")
    print(ctx["text"][:2000] + "\n…")
